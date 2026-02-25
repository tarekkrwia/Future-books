import streamlit as st
import google.generativeai as genai
import fitz  # PyMuPDF
import io
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Pt as PptPt
import json
import re

# --- إعدادات الصفحة ---
st.set_page_config(
    page_title="EduParser Pro - المُنسق الذكي",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- CSS لتحسين دعم العربية ---
st.markdown("""
<style>
    .stTextInput, .stTextArea, .stSelectbox { text-align: right; direction: rtl; }
    .stMarkdown { text-align: right; direction: rtl; }
    div[data-testid="stExpander"] details summary p { direction: rtl; }
</style>
""", unsafe_allow_html=True)

# --- إدارة حالة الجلسة (Session State) ---
if 'step' not in st.session_state:
    st.session_state.step = 1
if 'raw_text' not in st.session_state:
    st.session_state.raw_text = ""
if 'structured_data' not in st.session_state:
    st.session_state.structured_data = []

# --- دوال مساعدة (Helper Functions) ---
def extract_text_from_pdf(file):
    """استخراج النص من ملف PDF"""
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text

def clean_json_text(text):
    """تنظيف رد الذكاء الاصطناعي لاستخراج JSON صالح"""
    pattern = r"```json(.*?)```"
    match = re.search(pattern, text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return text.strip()

def create_word_doc(data):
    """إنشاء ملف Word من الأسئلة"""
    doc = Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(14)
    
    doc.add_heading('بنك الأسئلة - EduParser', 0)
    
    for i, item in enumerate(data, 1):
        p = doc.add_paragraph()
        runner = p.add_run(f"س{i}: {item.get('question', '')}")
        runner.bold = True
        runner.font.size = Pt(16)
        
        options = item.get('options', [])
        if options:
            for opt in options:
                doc.add_paragraph(f"- {opt}", style='List Bullet')
        
        ans = doc.add_paragraph()
        run_ans = ans.add_run(f"الإجابة الصحيحة: {item.get('answer', 'غير محدد')}")
        run_ans.font.color.rgb = RGBColor(0, 128, 0)
        
        doc.add_paragraph("-" * 50)
        
    bio = io.BytesIO()
    doc.save(bio)
    return bio

def create_ppt_pres(data):
    """إنشاء ملف PowerPoint"""
    prs = Presentation()
    for item in data:
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = item.get('question', 'سؤال')
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        options = item.get('options', [])
        for opt in options:
            p = tf.add_paragraph()
            p.text = opt
            p.font.size = PptPt(24)
            
        p_ans = tf.add_paragraph()
        p_ans.text = f"\n✅ الإجابة: {item.get('answer', '')}"
        p_ans.font.bold = True
        p_ans.font.color.rgb = RGBColor(0, 150, 0)
        
    bio = io.BytesIO()
    prs.save(bio)
    return bio

# --- الشريط الجانبي (Sidebar) ---
with st.sidebar:
    st.title("⚙️ الإعدادات")
    api_key = st.text_input("أدخل مفتاح Gemini API", type="password")
    st.info("احصل على المفتاح من: aistudio.google.com")
    
    st.markdown("---")
    st.markdown("### خطوات العمل:")
    steps = ["1. 📂 رفع الملفات", "2. 📝 مراجعة النص", "3. 🤖 الهيكلة بالذكاء", "4. 💾 التصدير"]
    current_step_idx = st.session_state.step - 1
    st.radio("المرحلة الحالية:", steps, index=current_step_idx, disabled=True)
    
    if st.button("🔄 إعادة البدء"):
        st.session_state.step = 1
        st.session_state.raw_text = ""
        st.session_state.structured_data = []
        st.rerun()

# --- المحتوى الرئيسي ---

if st.session_state.step == 1:
    st.header("1️⃣ رفع الملفات (PDF أو نصوص)")
    st.write("ارفع مذكراتك وسنقوم باستخراج النص منها.")
    
    uploaded_files = st.file_uploader("اختر الملفات", type=['pdf', 'txt'], accept_multiple_files=True)
    
    if uploaded_files:
        if st.button("بدء المعالجة واستخراج النص ➡️"):
            combined_text = ""
            progress_bar = st.progress(0)
            
            for idx, file in enumerate(uploaded_files):
                if file.type == "application/pdf":
                    combined_text += extract_text_from_pdf(file)
                elif file.type == "text/plain":
                    combined_text += str(file.read(), "utf-8")
                
                progress_bar.progress((idx + 1) / len(uploaded_files))
            
            st.session_state.raw_text = combined_text
            st.session_state.step = 2
            st.rerun()

elif st.session_state.step == 2:
    st.header("2️⃣ مراجعة النص الخام")
    st.write("قم بتعديل أي أخطاء ظهرت أثناء القراءة قبل إرسال النص للذكاء الاصطناعي.")
    
    edited_text = st.text_area("النص المستخرج:", value=st.session_state.raw_text, height=400)
    st.session_state.raw_text = edited_text
    
    col1, col2 = st.columns([1, 4])
    with col1:
        st.button("⬅️ رجوع", on_click=lambda: st.session_state.update(step=1))
    with col2:
        if st.button("تحليل وهيكلة النص (AI) ➡️"):
            if not api_key:
                st.error("الرجاء إدخال مفتاح API في القائمة الجانبية أولاً!")
            else:
                try:
                    genai.configure(api_key=api_key)
                    model = genai.GenerativeModel('gemini-1.5-flash')
                    
                    with st.spinner("جاري تحليل الأسئلة وتنسيقها... هذا قد يستغرق لحظات"):
                        prompt = f"""
                        أنت مساعد تعليمي خبير. قم بتحليل النص التالي واستخرج منه الأسئلة.
                        النص:
                        {edited_text}
                        
                        المطلوب:
                        قم بإخراج النتيجة بتنسيق JSON فقط (array of objects) بدون أي نصوص إضافية.
                        كل كائن يجب أن يحتوي على:
                        - "question": نص السؤال.
                        - "options": قائمة بالخيارات (إذا كان سؤال اختياري، وإلا اتركها قائمة فارغة).
                        - "answer": الإجابة الصحيحة (إذا كانت مذكورة في النص، وإلا اكتب "يحتاج مراجعة").
                        - "type": نوع السؤال ("mcq" أو "essay").
                        """
                        
                        response = model.generate_content(prompt)
                        json_str = clean_json_text(response.text)
                        st.session_state.structured_data = json.loads(json_str)
                        st.session_state.step = 3
                        st.rerun()
                except Exception as e:
                    st.error(f"حدث خطأ أثناء المعالجة: {e}")

elif st.session_state.step == 3:
    st.header("3️⃣ تنظيم وهيكلة الأسئلة")
    st.write("هنا تظهر الأسئلة بعد أن فهمها الذكاء الاصطناعي. يمكنك التعديل والحذف.")
    
    questions = st.session_state.structured_data
    
    for i, q in enumerate(questions):
        with st.expander(f"سؤال {i+1}: {q.get('question', '')[:50]}...", expanded=False):
            col_a, col_b = st.columns([3, 1])
            with col_a:
                q['question'] = st.text_input(f"نص السؤال {i+1}", q.get('question', ''))
                q['answer'] = st.text_input(f"الإجابة {i+1}", q.get('answer', ''))
            with col_b:
                q['type'] = st.selectbox(f"النوع {i+1}", ["mcq", "essay"], index=0 if q.get('type')=='mcq' else 1)
            
            if q['type'] == 'mcq':
                opts_str = "\n".join(q.get('options', []))
                new_opts = st.text_area(f"الخيارات (كل خيار في سطر) {i+1}", opts_str)
                q['options'] = new_opts.split('\n')
                
            if st.button(f"🗑️ حذف السؤال {i+1}", key=f"del_{i}"):
                questions.pop(i)
                st.session_state.structured_data = questions
                st.rerun()

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 2, 2])
    with col1:
        st.button("⬅️ رجوع", on_click=lambda: st.session_state.update(step=2))
    with col2:
        if st.button("➕ إضافة سؤال يدوي"):
            questions.append({"question": "سؤال جديد", "options": [], "answer": "", "type": "essay"})
            st.session_state.structured_data = questions
            st.rerun()
    with col3:
        st.button("اعتماد والذهاب للتصدير ✅", on_click=lambda: st.session_state.update(step=4))

elif st.session_state.step == 4:
    st.header("4️⃣ التصدير والتحميل")
    st.success("تم تجهيز بنك الأسئلة بنجاح! اختر الصيغة المناسبة للتحميل.")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### 📄 ملف Word")
        st.write("مناسب للطباعة وتوزيع المذكرات.")
        docx_file = create_word_doc(st.session_state.structured_data)
        st.download_button(
            label="تحميل بصيغة Word (.docx)",
            data=docx_file.getvalue(),
            file_name="Question_Bank.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
    with col2:
        st.markdown("### 📊 عرض PowerPoint")
        st.write("مناسب للعرض في الحصة أو السنتر.")
        ppt_file = create_ppt_pres(st.session_state.structured_data)
        st.download_button(
            label="تحميل بصيغة PowerPoint (.pptx)",
            data=ppt_file.getvalue(),
            file_name="Lesson_Slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    st.markdown("---")
    st.button("🔄 البدء بملف جديد", on_click=lambda: st.session_state.update(step=1, raw_text="", structured_data=[]))


